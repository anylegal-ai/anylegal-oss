"""
Document Management Tool Implementations

Workspace file CRUD with dual-mode support for text files and DOCX blobs:
- list_documents
- read_document   (branching: DOCX text/xml view vs HTML as-is)
- create_document (workspace text files only — anylegal.md, Playbook/*.md)
- edit_document   (branching: text→tracked-change for DOCX vs str.replace for HTML)
- clone_document  (versioned clone for editing)
- create_folder, delete_document, delete_folder

DOCX creation uses ``run_code`` (python-docx / docx-js), not this file.

DOCX editing approach (edit_document):
  LLM sends plain text old_text / new_text. Backend:
  1. Finds old_text inside <w:t> elements (quote/case normalization)
  2. Generates OOXML w:del/w:ins tracked-change markup
  3. Preserves original formatting (w:rPr)
  4. Validates, repacks blob, saves session

For structural OOXML edits, see the docx-xml skill (run_code + lxml + zipfile).
"""

import io
import os
import re
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List

from ..session import WorkspaceSession

logger = logging.getLogger(__name__)

def extract_xlsx_text(blob: bytes, filename: str = "", max_rows: int = 500) -> str:
    """Extract spreadsheet content as markdown tables. Used at upload and on legacy re-read."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(blob), data_only=False, read_only=True)
        lines = []
        total_rows = 0
        for sheet in wb.sheetnames:
            if total_rows >= max_rows:
                lines.append(f"\n*[Truncated at {max_rows} rows — use run_python with openpyxl for full data]*")
                break
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(row)
                total_rows += 1
                if total_rows >= max_rows:
                    break
            if not rows:
                continue
            lines.append(f"## Sheet: {sheet}")

            header = [str(c) if c is not None else "" for c in rows[0]]
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in rows[1:]:
                cells = [str(c) if c is not None else "" for c in row]

                while len(cells) < len(header):
                    cells.append("")
                cells = cells[:len(header)]
                lines.append("| " + " | ".join(cells) + " |")
        wb.close()
        return "\n".join(lines) if lines else "[Empty spreadsheet]"
    except Exception as e:
        logger.warning(f"XLSX text extraction failed for {filename}: {e}")
        return f"[XLSX extraction failed: {e}]"

def extract_pptx_text(blob: bytes, filename: str = "") -> str:
    """Extract presentation content as structured markdown. Used at upload and on legacy re-read."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(blob))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            lines.append(f"## Slide {i}" + (f": {title}" if title else ""))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and text != title:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append("| " + " | ".join(cells) + " |")
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"\n> Notes: {notes}")
            lines.append("")
        return "\n".join(lines) if lines else "[Empty presentation]"
    except Exception as e:
        logger.warning(f"PPTX text extraction failed for {filename}: {e}")
        return f"[PPTX extraction failed: {e}]"

_DOC_MIME_TYPES = {"application/msword", "application/x-ole-storage", "application/octet-stream"}

def _ensure_docx_blob(doc, session: WorkspaceSession) -> bool:
    """
    Ensure a document has a valid docx_blob.

    If the document only has a binary_blob (legacy .doc upload), attempts
    to convert it to .docx via the LibreOffice service. On success, stores
    the converted bytes as docx_blob and saves the session.

    Returns True if docx_blob is available (existing or newly converted).
    """
    if doc.docx_blob is not None:
        return True

    blob = getattr(doc, "binary_blob", None)
    if not blob:
        return False

    mime = getattr(doc, "mime_type", None) or ""
    is_doc_path = any(
        getattr(doc, "description", "").lower().endswith(ext)
        for ext in (".doc", ".dot")
    )
    if mime not in _DOC_MIME_TYPES and not is_doc_path:
        return False

    import requests as http_requests
    libreoffice_url = os.environ.get("LIBREOFFICE_SERVICE_URL", "http://localhost:8002")
    filename = getattr(doc, "description", None) or "document.doc"
    try:
        resp = http_requests.post(
            f"{libreoffice_url}/convert",
            files={"file": (filename, blob, "application/msword")},
            params={"format": "docx"},
            timeout=120,
        )
        if resp.status_code == 200 and resp.headers.get("content-type", "").startswith("application/"):
            docx_bytes = resp.content

            html_content = doc.content
            try:
                from ..docx_service import DocxService
                html_content, _ = DocxService.docx_to_html(docx_bytes)
            except Exception:
                pass
            doc.update_docx(docx_bytes, html_content)
            doc.content = html_content
            session.save()
            logger.info(f".doc → .docx on-demand conversion succeeded for {filename}")
            return True
        else:
            logger.warning(f".doc → .docx conversion failed (status {resp.status_code})")
    except Exception as e:
        logger.warning(f".doc → .docx conversion unavailable: {e}")

    return False

_XML_EDITING_REFERENCE = """
## XML Editing Quick Reference (Advanced)

NOTE: For normal editing, just use edit_document with plain text.
The system generates tracked changes automatically.

This reference is only needed for direct XML manipulation (advanced).
""".strip()

def list_documents(session: WorkspaceSession, folder: str = None, **kwargs) -> Dict[str, Any]:
    """List all documents in the workspace, optionally filtered by folder."""
    docs = []
    for path, doc in session.documents.items():

        if folder:
            prefix = folder.rstrip('/') + '/'
            if not path.startswith(prefix):
                continue

        parts = path.replace("\\", "/").split("/")
        doc_folder = "/".join(parts[:-1]) + "/" if len(parts) > 1 else ""

        docs.append({
            "path": path,
            "folder": doc_folder,
            "description": doc.description,
            "created_at": doc.created_at.isoformat(),
            "modified_at": doc.modified_at.isoformat(),
            "size": len(doc.content),
            "is_active": path == session.active_document,
            "format": doc.format,
            "has_docx": doc.docx_blob is not None,
            "has_binary": doc.binary_blob is not None,
            "mime_type": doc.mime_type,
            "is_synced": doc.is_synced
        })

    workspace_files = []
    for wf_path, wf_content in session.workspace_files.items():
        if folder:
            prefix = folder.rstrip('/') + '/'
            if not wf_path.startswith(prefix) and wf_path != folder.rstrip('/'):
                continue
        workspace_files.append({
            "path": wf_path,
            "type": "workspace_file",
            "size": len(wf_content),
            "editable": True,
        })

    skills = session.get_skill_files()

    templates = session.get_template_files()

    return {
        "success": True,
        "documents": docs,
        "workspace_files": workspace_files,
        "skills": skills,
        "templates": templates,
        "count": len(docs),
        "workspace_file_count": len(workspace_files),
        "active_document": session.active_document,
        "folders": sorted(session.folders),
    }

def _apply_range(
    content: str,
    around_text: Optional[str] = None,
    context_chars: int = 2000,
    start_text: Optional[str] = None,
    end_text: Optional[str] = None,
    paragraph_range: Optional[List[int]] = None,
) -> Dict[str, Any]:
    """Slice ``content`` according to the requested range parameters.

    Returns ``{content, range_info}``. If a requested range can't be
    satisfied (anchor not found, indices out of bounds), returns
    ``{error: str}`` instead. Caller should wrap in the standard
    ``{success: False, error}`` shape.

    Modes (mutually exclusive — caller must validate that at most one
    is set; if multiple are set, ``around_text`` wins, then range, then
    paragraphs):
    - around_text: window of ``context_chars`` total around the first
      occurrence of ``around_text``.
    - start_text / end_text: content between two anchors (inclusive). If
      ``end_text`` is missing, return from ``start_text`` to EOF.
    - paragraph_range: ``[start_idx, end_idx]`` slice of paragraphs split
      on single ``\\n``. Inclusive range.
    """
    total = len(content)

    if around_text:
        idx = content.find(around_text)
        if idx < 0:
            return {"error": f"around_text not found in document: {around_text!r}"}
        match_end = idx + len(around_text)
        half = max(0, context_chars - len(around_text)) // 2
        start = max(0, idx - half)
        end = min(total, match_end + half)
        return {
            "content": content[start:end],
            "range_info": {
                "mode": "around_text",
                "anchor": around_text,
                "anchor_offset": idx,
                "start_offset": start,
                "end_offset": end,
                "total_size": total,
            },
        }

    if start_text is not None:
        s_idx = content.find(start_text)
        if s_idx < 0:
            return {"error": f"start_text not found in document: {start_text!r}"}
        if end_text:
            e_idx = content.find(end_text, s_idx + len(start_text))
            if e_idx < 0:

                end = total
                end_anchor_found = False
            else:
                end = e_idx + len(end_text)
                end_anchor_found = True
        else:
            end = total
            end_anchor_found = None
        return {
            "content": content[s_idx:end],
            "range_info": {
                "mode": "start_end_text",
                "start_anchor": start_text,
                "end_anchor": end_text,
                "end_anchor_found": end_anchor_found,
                "start_offset": s_idx,
                "end_offset": end,
                "total_size": total,
            },
        }

    if paragraph_range is not None:
        if len(paragraph_range) != 2:
            return {"error": "paragraph_range must be [start_idx, end_idx]"}
        p_start, p_end = paragraph_range
        paragraphs = content.split("\n")
        n = len(paragraphs)
        if p_start < 0 or p_start >= n:
            return {"error": f"paragraph_range start {p_start} out of bounds (0..{n-1})"}

        p_end_clamped = min(p_end, n - 1)
        sliced = paragraphs[p_start : p_end_clamped + 1]
        return {
            "content": "\n".join(sliced),
            "range_info": {
                "mode": "paragraph_range",
                "requested": [p_start, p_end],
                "applied": [p_start, p_end_clamped],
                "paragraph_count_total": n,
                "paragraph_count_returned": len(sliced),
            },
        }

    return {"content": content, "range_info": None}

def read_document(
    session: WorkspaceSession,
    path: str,
    view: str = "text",
    around_text: Optional[str] = None,
    context_chars: int = 2000,
    start_text: Optional[str] = None,
    end_text: Optional[str] = None,
    paragraph_range: Optional[List[int]] = None,
    **kwargs
) -> Dict[str, Any]:
    """
    Read the content of a document.

    Dual-mode branching:
    - DOCX-native (has docx_blob):
        - view="text" (default): plain text for analysis
        - view="xml": pretty-printed document.xml for editing
    - HTML-native (no docx_blob): return HTML content as-is

    Range params (optional, mutually exclusive — at most one should be
    set; multiple set evaluates in this priority order: around_text,
    start_text/end_text, paragraph_range):
    - around_text + context_chars: slice ±context_chars/2 around the
      first occurrence of around_text. Use after edits to verify a
      specific region without re-reading the whole document.
    - start_text + end_text: slice between two anchors, inclusive. If
      end_text is missing in the doc, returns start_text to EOF.
    - paragraph_range: [start_idx, end_idx] slice of paragraphs (split
      on \\n), inclusive. Out-of-bounds end is clamped to last paragraph.

    Range params apply only to text view of DOCX, workspace_files, and
    HTML docs. They're ignored for view="xml" (where you want the full
    OOXML for structural work).
    """

    if not path or not path.strip() or len(path.strip()) < 3:
        if session.active_document:
            logger.info(f"[READ-DOC] No valid path provided, using active document: {session.active_document}")
            path = session.active_document

    range_kwargs = dict(
        around_text=around_text,
        context_chars=context_chars,
        start_text=start_text,
        end_text=end_text,
        paragraph_range=paragraph_range,
    )
    range_active = bool(around_text or start_text or paragraph_range)

    def _maybe_slice(resp: Dict[str, Any]) -> Dict[str, Any]:
        """Apply range params to a successful text response, if requested."""
        if not range_active or not resp.get("success"):
            return resp
        sliced = _apply_range(resp.get("content", ""), **range_kwargs)
        if "error" in sliced:
            return {"success": False, "error": sliced["error"], "path": resp.get("path")}
        resp["content"] = sliced["content"]
        resp["size"] = len(sliced["content"])
        resp["range_info"] = sliced["range_info"]
        return resp

    doc = session.get_document(path)

    if not doc:

        wf_content = session.get_workspace_file(path)
        if wf_content is not None:
            return _maybe_slice({
                "success": True,
                "path": path,
                "content": wf_content,
                "size": len(wf_content),
                "doc_type": "workspace_file",
                "view": "text",
                "editable": True,
            })

        _canonical_path = path if path.startswith("Skills/") else (
            "Skills/" + path[len("skills/"):] if path.startswith("skills/") else None
        )
        if _canonical_path:
            content = session.read_skill_file(_canonical_path)
            if content is not None:
                return _maybe_slice({
                    "success": True,
                    "path": path,
                    "content": content,
                    "size": len(content),
                    "doc_type": "skill",
                    "view": "text",
                    "editable": False,
                })

        return {
            "success": False,
            "error": f"Document not found: {path}",
            "available_documents": list(session.documents.keys()),
            "available_workspace_files": list(session.workspace_files.keys()),
        }

    _ensure_docx_blob(doc, session)

    if doc.docx_blob is not None:
        doc_type = "docx"

        if view == "xml":

            xml_content = doc.document_xml
            if xml_content is None:
                return {
                    "success": False,
                    "error": "Failed to extract XML from DOCX blob",
                    "doc_type": "docx"
                }

            return {
                "success": True,
                "path": path,
                "content": xml_content,
                "xml_editing_reference": _XML_EDITING_REFERENCE,
                "description": doc.description,
                "modified_at": doc.modified_at.isoformat(),
                "size": len(xml_content),
                "doc_type": "docx",
                "view": "xml",
                "format": doc.format,
                "has_docx": True,
                "is_synced": doc.is_synced,
                "docx_size": len(doc.docx_blob) if doc.docx_blob else 0,
                "hint": (
                    "This is the raw XML of word/document.xml (advanced view). "
                    "For editing, use edit_document with plain text old_text and "
                    "new_text — the system generates tracked changes automatically."
                ),
            }

        try:
            from ..docx_xml_service import extract_plain_text
            content = extract_plain_text(doc.docx_blob)
        except Exception as e:
            logger.warning(f"DOCX text extraction failed for {path}: {e}")
            content = doc.content

        return _maybe_slice({
            "success": True,
            "path": path,
            "content": content,
            "description": doc.description,
            "modified_at": doc.modified_at.isoformat(),
            "size": len(content),
            "doc_type": doc_type,
            "view": "text",
            "format": doc.format,
            "has_docx": True,
            "is_synced": doc.is_synced,
            "docx_size": len(doc.docx_blob) if doc.docx_blob else 0,
            "hint": (
                "This is plain text extracted from the DOCX. "
                "To edit, use edit_document with old_text (exact text to change) "
                "and new_text (replacement). Tracked changes are generated automatically."
            ),
        })

    if doc.binary_blob and doc.content.startswith("[Binary file:"):
        if doc.format == "xlsx" or (doc.mime_type and "spreadsheet" in doc.mime_type):
            doc.content = extract_xlsx_text(doc.binary_blob, doc.description or path)
        elif doc.format == "pptx" or (doc.mime_type and "presentation" in doc.mime_type):
            doc.content = extract_pptx_text(doc.binary_blob, doc.description or path)

    return _maybe_slice({
        "success": True,
        "path": path,
        "content": doc.content,
        "description": doc.description,
        "modified_at": doc.modified_at.isoformat(),
        "size": len(doc.content),
        "doc_type": "html",
        "view": "text",
        "format": doc.format,
        "has_docx": False,
        "is_synced": doc.is_synced,
    })

def create_document(
    session: WorkspaceSession,
    path: str,
    content: str,
    description: str = "",
    **kwargs
) -> Dict[str, Any]:
    """
    Create or overwrite a workspace text file.

    Primary use: anylegal.md, Playbook/*.md, and other plain-text workspace
    files. For DOCX documents, use ``run_code`` (python-docx / docx-js) or
    ``clone_document`` — this tool does NOT produce DOCX blobs.

    (Historically named ``write_document``; renamed Feb 2026. The alias was
    removed Apr 2026 once all callers and model calls migrated.)
    """
    BINARY_DOC_EXTS = (".docx", ".xlsx", ".pptx", ".pdf")
    if any(path.lower().endswith(ext) for ext in BINARY_DOC_EXTS):
        ext = path.lower().rsplit(".", 1)[-1]
        return {
            "success": False,
            "error": (
                f"create_document does not produce binary documents. "
                f"For .{ext} files, use run_code with python-docx / openpyxl / "
                f"python-pptx, or call clone_document to copy a template. "
                f"For text/markdown content, save to a .md path instead."
            ),
        }

    is_workspace_path = (
        path == "anylegal.md"
        or path.endswith("/anylegal.md")
        or path.startswith("Playbook/")
    )
    if is_workspace_path:
        already_exists = path in session.workspace_files
        session.set_workspace_file(path, content)
        if path == "Playbook/positions.md":
            session.set_playbook(content)

        parts = path.replace("\\", "/").split("/")
        if len(parts) > 1:
            folder = "/".join(parts[:-1]) + "/"
            session.folders.add(folder)
        session.save()
        result = {
            "success": True,
            "path": path,
            "action": "overwritten" if already_exists else "created",
            "size": len(content),
            "doc_type": "workspace_file",
        }

        if path.endswith("/anylegal.md") and path != "anylegal.md":
            root_instructions = session.workspace_files.get("anylegal.md", "")
            if not root_instructions or not root_instructions.strip():
                result["warning"] = (
                    f"You created folder-level instructions at '{path}', but the root "
                    f"anylegal.md (shown as 'Instructions' in the sidebar) is still empty. "
                    f"If this was meant to be the user's main instructions, use path "
                    f"'anylegal.md' instead. Folder-level instructions only apply to "
                    f"documents inside that specific folder."
                )
        return result

    if path.startswith("Skills/"):
        return {
            "success": False,
            "error": f"'{path}' is read-only. Skills/ contains system skill files.",
        }

    if path.startswith("Templates/"):
        filename = path.replace("\\", "/").split("/")[-1]
        return {
            "success": False,
            "error": (
                f"Cannot write to Templates/ — only users can manage templates. "
                f"Save the document to a regular path instead, e.g. '{filename}' or "
                f"'Client Projects/{filename}'."
            ),
        }

    is_markdown_output = path.lower().endswith('.md')
    is_new = path not in session.documents

    session.add_document(
        path=path,
        content=content,
        description=description,
        set_active=True
    )

    parts = path.replace("\\", "/").split("/")
    if len(parts) > 1:
        folder = "/".join(parts[:-1]) + "/"
        session.folders.add(folder)

    result = {
        "success": True,
        "path": path,
        "action": "created" if is_new else "overwritten",
        "size": len(content),
        "is_active": True,
    }

    if is_markdown_output:

        doc = session.get_document(path)
        if doc:
            doc.format = "markdown"
            doc.mime_type = "text/markdown"
            doc.binary_blob = content.encode('utf-8')
        result["format"] = "markdown"
        result["has_docx"] = False
    else:

        try:
            from ..docx_service import DocxService

            doc = session.get_document(path)
            if doc:
                docx_bytes = DocxService.markdown_to_docx(content)
                doc.update_docx(docx_bytes, content)
                result["format"] = "docx"
                result["docx_size"] = len(docx_bytes)
                result["has_docx"] = True
        except Exception as e:
            logger.error(f"Markdown→DOCX conversion failed for {path}: {e}", exc_info=True)
            result["format"] = "markdown"
            result["has_docx"] = False
            result["docx_error"] = str(e)

    session.save()
    return result

def edit_document(
    session: WorkspaceSession,
    path: str,
    old_text: str = "",
    new_text: str = "",
    explanation: str = "",
    start_text: str = "",
    end_text: str = "",
    near_text: str = "",
    **kwargs
) -> Dict[str, Any]:
    """
    Find and replace content in a document.

    Dual-mode branching:
    - DOCX-native (has docx_blob):
        1. Text-level edit: find old_text in <w:t> content, generate tracked
           changes (w:del/w:ins) automatically.  The LLM sends plain text.
        2. Raw XML fallback: if old_text is an XML fragment found in the raw
           document.xml, do direct string replacement (advanced).
    - HTML-native (no docx_blob): str.replace on HTML content.
    """

    if not path or not path.strip() or len(path.strip()) < 3:
        if session.active_document:
            logger.info(f"[DOCX-EDIT] No valid path provided, using active document: {session.active_document}")
            path = session.active_document
        else:
            return {"success": False, "error": "No document path provided and no active document set."}

    doc = session.get_document(path)
    if not doc:

        wf_content = session.get_workspace_file(path)
        if wf_content is not None:
            if old_text not in wf_content:
                return {
                    "success": False,
                    "error": f"Text not found in workspace file: {path}",
                    "doc_type": "workspace_file",
                }
            count = wf_content.count(old_text)
            if count > 1:
                return {
                    "success": False,
                    "error": f"Text appears {count} times. Include more context to make it unique.",
                    "doc_type": "workspace_file",
                }
            new_content = wf_content.replace(old_text, new_text, 1)
            session.set_workspace_file(path, new_content)

            if path == "Playbook/positions.md":
                session.playbook = new_content
            session.save()
            return {
                "success": True,
                "path": path,
                "doc_type": "workspace_file",
                "explanation": explanation or "",
            }

        if path.startswith("Skills/"):
            return {
                "success": False,
                "error": f"'{path}' is read-only and cannot be edited.",
            }

        if path.startswith("Templates/"):
            return {
                "success": False,
                "error": f"'{path}' is a template and cannot be edited by the agent. Use read_document to read it.",
            }

        return {
            "success": False,
            "error": f"Document not found: {path}",
            "available_documents": list(session.documents.keys()),
            "available_workspace_files": list(session.workspace_files.keys()),
        }

    _ensure_docx_blob(doc, session)

    if doc.docx_blob is not None:

        resolved_path, cloned_from = resolve_or_clone_to_v2(session, path)
        if resolved_path != path:
            path = resolved_path
            doc = session.get_document(path)
            if doc is None or doc.docx_blob is None:
                return {
                    "success": False,
                    "error": (
                        f"auto-clone-to-v2 produced path {resolved_path!r} "
                        f"but the document is missing or has no DOCX blob"
                    ),
                }
        try:
            from ..docx_xml_service import (
                apply_text_edit,
                validate_document_xml,
            )

            xml_content = doc.document_xml
            if xml_content is None:
                return {
                    "success": False,
                    "error": "Failed to extract XML from DOCX blob",
                    "doc_type": "docx"
                }

            new_xml = None
            edit_info = {}
            info = {}                                                

            if start_text and end_text:
                from ..docx_xml_service import apply_range_delete
                result_xml, info = apply_range_delete(
                    xml_content, start_text, end_text
                )
                if result_xml is not None:
                    new_xml = result_xml
                    edit_info = info
                    logger.info(
                        f"[DOCX-EDIT] Range delete on '{path}': "
                        f"{info.get('paragraphs_deleted', 0)} paragraphs"
                    )
                else:
                    return {
                        "success": False,
                        "error": info.get("error", "Range delete failed."),
                        "suggestion": info.get("suggestion", ""),
                        "doc_type": "docx",
                    }

            if new_xml is None and old_text:
                result_xml, info = apply_text_edit(
                    xml_content, old_text, new_text, near_text=near_text
                )
                if result_xml is not None:
                    new_xml = result_xml
                    edit_info = info
                    logger.info(
                        f"[DOCX-EDIT] Text-level edit on '{path}': "
                        f"matched '{info.get('matched_text', '')[:60]}'"
                    )

            if new_xml is None and old_text in xml_content:
                from ..docx_xml_service import (
                    _toc_paragraph_ranges,
                    _in_tracked_change,
                )

                toc_ranges = _toc_paragraph_ranges(xml_content)

                body_count = 0
                body_pos = None
                search_start = 0
                while True:
                    pos = xml_content.find(old_text, search_start)
                    if pos == -1:
                        break
                    if not _in_tracked_change(pos, toc_ranges):
                        body_count += 1
                        if body_pos is None:
                            body_pos = pos
                    search_start = pos + 1

                if body_count > 1:
                    return {
                        "success": False,
                        "error": (
                            f"XML fragment appears {body_count} times in document body. "
                            "Include more context to make it unique."
                        ),
                        "doc_type": "docx",
                    }
                elif body_count == 1:

                    new_xml = (
                        xml_content[:body_pos]
                        + new_text
                        + xml_content[body_pos + len(old_text) :]
                    )
                    edit_info = {"mode": "raw_xml"}
                    logger.info(f"[DOCX-EDIT] Raw XML edit on '{path}' (TOC-aware)")

            if new_xml is None:
                if not old_text and not (start_text and end_text):
                    return {
                        "success": False,
                        "error": "Provide old_text/new_text for editing, or start_text/end_text for range deletion.",
                        "doc_type": "docx",
                    }
                error_msg = info.get("error", "Text not found in document.")
                return {
                    "success": False,
                    "error": error_msg,
                    "suggestion": info.get("suggestion", ""),
                    "nearby_text": info.get("nearby_text", ""),
                    "doc_type": "docx",
                }

            errors = validate_document_xml(new_xml)
            if errors:
                logger.warning(f"[DOCX-EDIT] Validation warnings: {errors}")
                parse_errors = [
                    e for e in errors if e.startswith("XML parse error")
                ]
                if parse_errors:
                    return {
                        "success": False,
                        "error": (
                            "Edit produced invalid XML. "
                            f"Errors: {'; '.join(parse_errors)}"
                        ),
                        "doc_type": "docx",
                        "validation_errors": errors,
                    }

            doc.update_document_xml(new_xml)

            try:
                session.save()
            except Exception as save_err:
                logger.warning(f"[DOCX-EDIT] Session save failed: {save_err}")

            result = {
                "success": True,
                "path": path,
                "explanation": explanation or "",
                "doc_type": "docx",
                "has_docx": True,
                "docx_updated": True,
            }

            if cloned_from:
                result["cloned_from"] = cloned_from
                result["message"] = (
                    f"First edit on '{cloned_from}' — created working copy "
                    f"'{path}'. Original preserved; further edits in this "
                    f"session will target the working copy."
                )
            if edit_info.get("matched_text"):
                result["matched_text"] = edit_info["matched_text"]
            if edit_info.get("revision_ids"):
                result["revision_ids"] = edit_info["revision_ids"]
            if new_text:
                result["replacement_text"] = new_text
            if errors:
                result["validation_warnings"] = errors

            try:
                from ..docx_xml_service import _context_around_revision
                ctx = _context_around_revision(
                    new_xml, edit_info.get("revision_ids", [])
                )
                if ctx:
                    result["context_around_edit"] = ctx
            except Exception:
                pass                                     

            return result

        except Exception as e:
            logger.error(f"DOCX edit failed for {path}: {e}", exc_info=True)
            return {
                "success": False,
                "error": f"DOCX edit failed: {str(e)}",
                "doc_type": "docx"
            }

    result = session.edit_document(path, old_text, new_text)

    if result["success"]:
        if explanation:
            result["explanation"] = explanation
        result["doc_type"] = "html"

        updated_doc = session.get_document(path)
        if updated_doc:
            result["content"] = updated_doc.content

        if doc and doc.docx_blob:
            result["has_docx"] = True
            result["is_synced"] = doc.is_synced

    return result

def clone_document(
    session: "WorkspaceSession",
    source_path: str,
    target_path: str = "",
    **kwargs,
) -> Dict[str, Any]:
    """
    Clone a document to create the next version in a version chain.

    Law firm versioning: original → v2 → v3 → v4. All versions preserved.
    The model always passes the ORIGINAL document path. The backend automatically:
    1. Finds the latest existing version (v2, v3, etc.)
    2. Clones FROM the latest version TO the next version number
    3. Sets the new version as the active document

    Example chain:
      clone_document("Contract.docx")  → clones original → Contract_v2.docx
      clone_document("Contract.docx")  → clones v2      → Contract_v3.docx
      clone_document("Contract.docx")  → clones v3      → Contract_v4.docx
    """
    logger.info(f"clone_document: source_path={source_path!r}, target_path={target_path!r}")

    base_path = _strip_version_suffix(source_path)

    base_doc = session.get_document(base_path)
    if not base_doc:

        base_doc = session.get_document(source_path)
        if base_doc:
            base_path = source_path
        else:

            latest_path, next_version = _find_latest_version(session, base_path)
            if latest_path != base_path and session.get_document(latest_path):

                logger.info(f"clone_document: base {base_path!r} not found, but found {latest_path}")
            elif session.active_document:

                logger.warning(f"clone_document: garbled path {source_path!r}, falling back to active document: {session.active_document}")
                source_path = session.active_document
                base_path = _strip_version_suffix(source_path)
                base_doc = session.get_document(base_path)
                if not base_doc:
                    base_doc = session.get_document(source_path)
                    if base_doc:
                        base_path = source_path
            else:
                available = list(session.documents.keys())
                logger.warning(f"clone_document: source not found: {source_path}. Available: {available}")
                return {
                    "success": False,
                    "error": f"Source document not found: {source_path}. Available documents: {available}",
                }

    latest_path, next_version = _find_latest_version(session, base_path)
    latest_doc = session.get_document(latest_path)
    if not latest_doc:

        latest_doc = base_doc
        latest_path = base_path

    _ensure_docx_blob(latest_doc, session)

    if not latest_doc.docx_blob and not latest_doc.content and not latest_doc.binary_blob:
        return {
            "success": False,
            "error": f"Source document '{latest_path}' has no content to clone.",
        }

    if not target_path or not target_path.strip():
        target_path = _version_path(base_path, next_version)

    while target_path in session.documents:
        next_version += 1
        target_path = _version_path(base_path, next_version)

    session.add_document(
        path=target_path,
        content=latest_doc.content or "",
        description=f"v{next_version} of {Path(base_path).name}",
        set_active=True,
    )

    doc = session.get_document(target_path)
    if doc:
        if latest_doc.docx_blob:
            doc.update_docx(latest_doc.docx_blob, latest_doc.content or "")
            doc.format = "docx"
        elif latest_doc.binary_blob:
            doc.binary_blob = latest_doc.binary_blob
            doc.mime_type = getattr(latest_doc, "mime_type", None) or "application/octet-stream"
            doc.format = getattr(latest_doc, "format", "other")

    norm = target_path.replace("\\", "/")
    if "/" in norm:
        folder = norm.rsplit("/", 1)[0] + "/"
        session.folders.add(folder)

    session.save()

    cloned_from = latest_path if latest_path != base_path else "original"
    logger.info(f"clone_document: {cloned_from} → {target_path!r} (v{next_version}, docx={bool(latest_doc.docx_blob)})")
    return {
        "success": True,
        "path": target_path,
        "source_path": latest_path,
        "version": next_version,
        "format": "docx" if latest_doc.docx_blob else "other",
        "has_docx": bool(latest_doc.docx_blob),
    }

def _has_version_suffix(path: str) -> bool:
    """True if path ends in ``_v{N}.{ext}`` — model is working on an explicit version."""
    p = Path(path)
    return bool(re.search(r'_v\d+$', p.stem))

def resolve_or_clone_to_v2(
    session: "WorkspaceSession",
    path: str,
) -> tuple:
    """
    Map a mutating edit's target path to the actual write target, auto-cloning
    originals to ``_v2.docx`` on first edit. Protects the original against
    model-introduced corruption (bad OOXML, wrong span, hallucinated content)
    by preserving it as the pristine reference.

    Rules:
      * Path is ALREADY versioned (``Contract_v2.docx``, ``Contract_v3.docx``)
        → return as-is. The caller explicitly picked this version; we don't
        second-guess. Subsequent edits chain onto the same version until the
        user asks for a new one via ``clone_document`` (non-agent paths).
      * Path is the original (``Contract.docx``) AND a ``_v{N}`` already
        exists in this workspace:
          - If the latest ``_v{N}`` was finalized via accept_all_changes /
            reject_all_changes (no output_path) and not yet edited again,
            **clone to ``_v{N+1}`` and return the new path**. This is the
            Option B "round-bump" — each finalize event closes a round,
            the next edit starts a fresh version chain. Matches the
            lawyer's mental model of "version = round of negotiation."
          - Otherwise redirect to the existing ``_v{N}`` — edits within a
            round keep piling onto the same working copy.
      * Path is the original, no ``_v{N}`` exists, and the document has a
        DOCX blob → clone to ``_v2.docx`` and return the new path.
      * Path has no DOCX blob (missing doc, plain-text workspace file like
        ``anylegal.md``) → return path unchanged. Downstream tool handles
        the missing-blob error and no safety concern applies to text files.

    Returns ``(target_path, cloned_from)``. ``cloned_from`` is the source path
    when a new version was just created; ``None`` otherwise (including when
    the caller was redirected to an existing ``_v{N}``). Tools return both
    fields in their result so the UI can render a "v2 created" affordance
    distinct from "you're still editing the working copy."

    Disable via ``AUTO_CLONE_ON_FIRST_EDIT=false`` — useful for tests that
    want deterministic in-place edits.
    """
    import os

    if os.getenv("AUTO_CLONE_ON_FIRST_EDIT", "true").lower() == "false":
        return path, None

    if _has_version_suffix(path):
        return path, None

    base_path = path
    doc = session.get_document(base_path)
    if doc is None:
        return path, None                                     
    if not getattr(doc, "docx_blob", None):
        return path, None                                               

    latest_path, _next_version = _find_latest_version(session, base_path)
    if latest_path != base_path:

        latest_doc = session.get_document(latest_path)
        if latest_doc is not None and getattr(latest_doc, "finalized_at", None) is not None:
            new_version_path = _version_path(base_path, _next_version)
            clone_result = clone_document(
                session=session,
                source_path=base_path,
                target_path=new_version_path,
            )
            if clone_result.get("success"):
                target = clone_result["path"]
                logger.info(
                    f"resolve_or_clone_to_v2: round-bump on {base_path!r} — "
                    f"{latest_path} was finalized, cloned to {target!r} for new round"
                )
                return target, latest_path

            logger.warning(
                f"resolve_or_clone_to_v2: round-bump clone failed "
                f"({clone_result.get('error')}); piling onto {latest_path}"
            )
        return latest_path, None

    result = clone_document(session=session, source_path=base_path)
    if not result.get("success"):
        logger.warning(
            f"resolve_or_clone_to_v2: clone of {base_path!r} failed "
            f"({result.get('error')}); editing in place"
        )
        return path, None

    target = result["path"]
    logger.info(
        f"resolve_or_clone_to_v2: first edit on {base_path!r} → cloned to "
        f"{target!r} (original preserved)"
    )
    return target, base_path

def _strip_version_suffix(path: str) -> str:
    """Remove _v2, _v3 etc. from a path: Contract_v3.docx → Contract.docx."""
    p = Path(path)
    stem = re.sub(r'_v\d+$', '', p.stem)
    parent = str(p.parent) if str(p.parent) != "." else ""
    name = f"{stem}{p.suffix}"
    return f"{parent}/{name}" if parent else name

def _version_path(base_path: str, version: int) -> str:
    """Generate a versioned path: Contract.docx + version=3 → Contract_v3.docx."""
    p = Path(base_path)
    stem = re.sub(r'_v\d+$', '', p.stem)                             
    suffix = p.suffix or ".docx"
    parent = str(p.parent) if str(p.parent) != "." else ""
    name = f"{stem}_v{version}{suffix}"
    return f"{parent}/{name}" if parent else name

def _find_latest_version(session: "WorkspaceSession", base_path: str) -> tuple:
    """
    Find the highest version of a document in the workspace.

    Returns (latest_path, next_version_number).
    If no versions exist, returns (base_path, 2).
    """
    p = Path(base_path)
    stem = re.sub(r'_v\d+$', '', p.stem)
    suffix = p.suffix or ".docx"
    parent = str(p.parent) if str(p.parent) != "." else ""

    highest_version = 1                 
    latest_path = base_path

    for doc_path in session.documents:
        doc_p = Path(doc_path)
        doc_parent = str(doc_p.parent) if str(doc_p.parent) != "." else ""
        if doc_parent != parent or doc_p.suffix.lower() != suffix.lower():
            continue
        match = re.match(rf'^{re.escape(stem)}_v(\d+)$', doc_p.stem)
        if match:
            v = int(match.group(1))
            if v > highest_version:
                highest_version = v
                latest_path = doc_path

    return latest_path, highest_version + 1

def create_folder(session: WorkspaceSession, folder_path: str, **kwargs) -> Dict[str, Any]:
    """Create a folder in the workspace. Used by the /setup skill to scaffold folder structure."""
    if not folder_path or not folder_path.strip():
        return {"success": False, "error": "folder_path is required"}

    clean = folder_path.replace("\\", "/").strip("/") + "/"

    top = clean.split("/")[0].lower()
    if top == "skills":
        return {"success": False, "error": "Cannot create folders inside Skills/ — it is a system folder"}

    session.create_folder(clean)
    session.save()
    return {
        "success": True,
        "folder_path": clean,
        "message": f"Folder '{clean}' created successfully",
    }

def delete_document(session: WorkspaceSession, path: str, **kwargs) -> Dict[str, Any]:
    """Delete a single document or workspace file from the workspace."""
    if not path or not path.strip():
        return {"success": False, "error": "path is required"}
    clean = path.replace("\\", "/").strip("/")

    if clean.split("/")[0] == "Skills":
        return {"success": False, "error": "Cannot delete files from the Skills/ system folder"}
    removed = session.remove_document(clean)
    if not removed:

        if clean in session.workspace_files:
            del session.workspace_files[clean]
            if clean == "anylegal.md" or clean == "agents.md":
                session.agents_md = ""
            removed = True
    if removed:
        return {"success": True, "path": clean, "message": f"'{clean}' deleted"}
    return {"success": False, "error": f"'{clean}' not found in workspace"}

def delete_folder(session: WorkspaceSession, folder_path: str, **kwargs) -> Dict[str, Any]:
    """Delete a user folder and all its contents. System folders (Skills/, Templates/) are protected."""
    if not folder_path or not folder_path.strip():
        return {"success": False, "error": "folder_path is required"}
    try:
        count = session.delete_folder(folder_path)
        clean = folder_path.replace("\\", "/").strip("/") + "/"
        return {
            "success": True,
            "folder_path": clean,
            "documents_deleted": count,
            "message": f"Folder '{clean}' and {count} file(s) deleted",
        }
    except ValueError as e:
        return {"success": False, "error": str(e)}

def instantiate_template(
    session: WorkspaceSession,
    template_path: str,
    output_path: str,
    replacements: Dict[str, str],
    **kwargs,
) -> Dict[str, Any]:
    """
    Create a new DOCX from a template by filling placeholders, with NO tracked
    changes in the output. The original template is untouched.

    Implementation: clone the template's bytes to a new doc at ``output_path``,
    apply ``apply_text_edit`` for each replacement (which generates valid
    tracked-change OOXML preserving run properties — ``<w:rPr>``), then accept
    all the freshly-created revision IDs to produce a clean final output.

    This reuses the battle-tested matcher from ``edit_document`` (multi-run
    safe, smart-quote tolerant, whitespace-agnostic) without duplicating
    ~200 lines of match-finding logic. The output has no ``<w:ins>``/``<w:del>``
    markup — the document looks identical to a fresh manual fill.

    Args:
        session: Workspace session
        template_path: Path to the source template (e.g. "Templates/Board_Resolution.docx")
        output_path: Path for the new document (e.g. "Acme Board Res 2026-04-25.docx")
        replacements: Dict of placeholder text → replacement text. Each key
            must appear in the template; missing keys are reported in
            ``not_found``.

    Returns:
        ``{success, output_path, applied: [...], not_found: [...], doc_type: "docx"}``
    """
    from ..docx_xml_service import (
        apply_text_edit,
        accept_specific_changes,
        repack_docx,
        validate_document_xml,
    )

    try:
        if not template_path:
            return {"success": False, "error": "template_path is required"}
        if not output_path:
            return {"success": False, "error": "output_path is required"}
        if not replacements:
            return {"success": False, "error": "replacements must be a non-empty dict"}

        template_doc = session.get_document(template_path)
        if not template_doc:
            return {
                "success": False,
                "error": f"Template not found: {template_path}",
                "available_documents": list(session.documents.keys()),
            }
        _ensure_docx_blob(template_doc, session)
        if not template_doc.docx_blob:
            return {
                "success": False,
                "error": f"Template {template_path!r} has no DOCX blob — cannot instantiate.",
            }

        original_blob = template_doc.docx_blob
        from ..docx_xml_service import extract_document_xml
        try:
            xml_content = extract_document_xml(original_blob)
        except Exception:
            xml_content = template_doc.document_xml
        if xml_content is None:
            return {
                "success": False,
                "error": "Failed to extract document.xml from template",
            }

        applied: List[str] = []
        not_found: List[str] = []
        all_revision_ids: List[int] = []

        for old_text, new_text in replacements.items():
            if not old_text:
                continue
            new_xml, info = apply_text_edit(xml_content, old_text, str(new_text))
            if new_xml is None:
                not_found.append(old_text)
                continue
            xml_content = new_xml
            applied.append(old_text)
            for rid in info.get("revision_ids", []):
                try:
                    all_revision_ids.append(int(rid))
                except (TypeError, ValueError):
                    pass

        if all_revision_ids:
            xml_content, _ = accept_specific_changes(xml_content, all_revision_ids)

        errors = validate_document_xml(xml_content)
        parse_errors = [e for e in errors if e.startswith("XML parse error")]
        if parse_errors:
            return {
                "success": False,
                "error": f"XML invalid after instantiation: {parse_errors[0]}",
                "validation_errors": errors,
                "applied": applied,
                "not_found": not_found,
            }

        new_blob = repack_docx(original_blob, xml_content)

        if output_path in session.documents:
            existing = session.get_document(output_path)
            existing.update_docx(new_blob)
        else:
            session.add_document(
                path=output_path,
                content="",
                description=f"Instantiated from {template_path}",
                set_active=True,
            )
            session.get_document(output_path).update_docx(new_blob)

        session.save()

        return {
            "success": True,
            "output_path": output_path,
            "doc_type": "docx",
            "applied": applied,
            "not_found": not_found,
            "template_path": template_path,
        }

    except Exception as e:
        logger.error(f"instantiate_template failed: {e}", exc_info=True)
        return {"success": False, "error": str(e)}

DOCUMENT_TOOLS = {
    "list_documents": list_documents,
    "read_document": read_document,
    "create_document": create_document,
    "edit_document": edit_document,
    "clone_document": clone_document,
    "create_folder": create_folder,
    "delete_document": delete_document,
    "delete_folder": delete_folder,
    "instantiate_template": instantiate_template,
}
